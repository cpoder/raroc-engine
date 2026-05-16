"""Build the OpenRAROC demo deck (PPTX).

Audience: corporate treasurers / CFOs and the partners (advisors, consultancies,
fractional CFO firms) who bring this tool to them. NOT bankers.
Length: 12 slides, 10–15 min talk track.
Brand: dark theme matching openraroc.com (slate background, cyan accent).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette (matches the OpenRAROC web app) ─────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
SURFACE  = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
SURFACE2 = RGBColor(0x33, 0x41, 0x55)   # slate-700
BORDER   = RGBColor(0x47, 0x55, 0x69)   # slate-600
TEXT     = RGBColor(0xE2, 0xE8, 0xF0)   # slate-200
TEXT2    = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
TEXT3    = RGBColor(0x64, 0x74, 0x8B)   # slate-500
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)   # sky-400
ACCENT2  = RGBColor(0x0E, 0xA5, 0xE9)   # sky-500
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
YELLOW   = RGBColor(0xEA, 0xB3, 0x08)
RED      = RGBColor(0xEF, 0x44, 0x44)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)


# ── Helpers ───────────────────────────────────────────────────────────
def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=SURFACE, line=BORDER, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_multiline(slide, x, y, w, h, lines, *, size=14, color=TEXT,
                  align=PP_ALIGN.LEFT, font="Calibri", line_spacing=1.2):
    """lines: list of (text, {bold, color, size}) or just str."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = opts.get("color", color)
    return tb


def add_chip(slide, x, y, text, color=ACCENT, bg=SURFACE, w=Inches(1.6), h=Inches(0.32)):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.adjustments[0] = 0.5
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.color.rgb = color; rect.line.width = Pt(0.75)
    tf = rect.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = color
    r.font.name = "Calibri"
    return rect


def add_header(slide, chapter, title, subtitle=None):
    """Top strip: chapter chip on left, title underneath."""
    if chapter:
        add_chip(slide, Inches(0.5), Inches(0.4), chapter,
                 color=ACCENT, bg=SURFACE, w=Inches(1.8), h=Inches(0.32))
    add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.4),
                 subtitle, size=14, color=TEXT2)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(7), Inches(0.3),
             "OpenRAROC  •  See your credit through your banker's eyes",
             size=10, color=TEXT3)
    add_text(slide, Inches(11), Inches(7.0), Inches(2), Inches(0.3),
             f"{page_num} / {total}", size=10, color=TEXT3, align=PP_ALIGN.RIGHT)


def add_line(slide, x1, y1, x2, y2, color=BORDER, width=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_table(slide, x, y, w, h, data, col_widths=None,
              header_bg=SURFACE2, header_color=ACCENT,
              row_bg=SURFACE, row_alt=None, text_color=TEXT,
              highlight_rows=None, highlight_color=GREEN,
              font_size=12, header_size=11, num_cols_right=None):
    """data: list of rows, first is header.
    highlight_rows: list of indices (excluding header) to highlight.
    num_cols_right: list of column indices to right-align.
    """
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    highlight_rows = highlight_rows or []
    num_cols_right = num_cols_right or []

    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Emu(45000); cell.margin_right = Emu(45000)
            cell.margin_top = Emu(30000); cell.margin_bottom = Emu(30000)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
                txt_color = header_color
                bold = True; size = header_size
            elif (ri - 1) in highlight_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = SURFACE2
                txt_color = highlight_color
                bold = True; size = font_size
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_alt if (row_alt and ri % 2 == 0) else row_bg
                txt_color = text_color
                bold = False; size = font_size

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if ci in num_cols_right else PP_ALIGN.LEFT
            p.text = ""
            r = p.add_run(); r.text = str(val)
            r.font.name = "Calibri"
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = txt_color
    return tbl


def add_kpi(slide, x, y, w, h, label, value, sub=None, value_color=WHITE):
    add_rect(slide, x, y, w, h, fill=SURFACE, line=BORDER)
    add_text(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.25),
             label.upper(), size=9, bold=True, color=TEXT3)
    add_text(slide, x + Inches(0.2), y + Inches(0.42), w - Inches(0.4), Inches(0.55),
             value, size=22, bold=True, color=value_color)
    if sub:
        add_text(slide, x + Inches(0.2), y + h - Inches(0.35),
                 w - Inches(0.4), Inches(0.25),
                 sub, size=9, color=TEXT2)


def add_bullet_block(slide, x, y, w, h, title, bullets,
                     title_color=ACCENT, bullet_color=TEXT, fill=SURFACE):
    add_rect(slide, x, y, w, h, fill=fill, line=BORDER)
    add_text(slide, x + Inches(0.25), y + Inches(0.2), w - Inches(0.5), Inches(0.4),
             title, size=14, bold=True, color=title_color)
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.7),
                                  w - Inches(0.5), h - Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        p.space_after = Pt(4)
        r = p.add_run(); r.text = "•  " + b
        r.font.size = Pt(12); r.font.color.rgb = bullet_color
        r.font.name = "Calibri"


# ── Build the deck ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 12
blank = prs.slide_layouts[6]


def new_slide(chapter=None, title=None, subtitle=None, page=None, notes=None):
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    if title:
        add_header(slide, chapter or "", title, subtitle)
    if page is not None:
        add_footer(slide, page, TOTAL)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# ── Speaker notes — corporate + partner audience ────────────────────
NOTES_S1 = """
TIMING: ~30s.

WHO IS IN THE ROOM:
  Two profiles, one deck. Left half — CFOs, treasurers, finance directors at mid-cap companies who borrow. Right half — advisors who serve them: treasury consultancies, debt advisory boutiques, fractional CFO firms, accountants with corporate finance practices. Both want the same outcome: better deals from banks.

OPENING LINE (memorize, don't paraphrase):
  "Your bank quotes you a number. Until today, you couldn't tell whether it was generous, fair, or daylight robbery. After this demo, you'll know — and so will your next quote."

THEN STOP. Let it land. Don't introduce yourself for at least 10 seconds.

DON'T:
  • Don't say "we're going to teach you about RAROC". Sounds patronising; half the audience already knows what it is. Just demonstrate.
  • Don't pitch features. We're 30 seconds in.
""".strip()

NOTES_S2 = """
TIMING: ~75s.

THE SETUP:
  Two columns side by side. Left = "what you see today" (a single number quoted by the RM). Right = "what your bank actually sees" (the whole machinery: capital, expected loss, cost-to-income, hurdle rate). Big asymmetry. That's the whole product opportunity.

VOCABULARY YOU'LL NEED FOR THIS SLIDE:
  • Spread = the margin the bank charges over a reference rate (EURIBOR, SOFR). 150 bp = 1.5%. The headline number on every term sheet.
  • RAROC = Risk-Adjusted Return on Capital. Bank's internal pricing metric. The full glossary is in GLOSSARY.md but at this slide you can just say "the post-tax return your bank earns on the capital it sets aside for your loan".
  • Pillar 3 / CR6 = annual public disclosure that every IRB-regulated bank must publish. Has the inputs we need to reverse-engineer their pricing.

LINE TO LAND ON (the bottom callout):
  "30 basis points on a EUR 50M facility is EUR 750k over the tenor. That's a hire. That's a software stack. That's the dividend your CEO didn't get to announce."

IF AN ADVISOR IN THE ROOM ASKS "DO TREASURERS REALLY NOT SEE THIS?":
  → "They see the price. They've never had the model behind it. We did the model."
""".strip()

NOTES_S3 = """
TIMING: ~60s. Move quickly through the four KPIs; spend more time on the For-You / For-Your-Firm split at the bottom.

VOCABULARY:
  • Basel III IRB calculator = the same regulatory math your bank uses internally; we just compute it from the customer side.
  • Min-spread solver = inverts the formula. "What's the lowest spread my bank can accept?". Brent's method, milliseconds per call.
  • MCP = lets AI assistants like Claude / ChatGPT call the engine directly. Niche today, mainstream in 12 months.
  • Pro tier (EUR 49/year) = unlocks all 59 banks; free tier covers the four largest (BNP, HSBC, Deutsche, JPM).

DELIVERY:
  Skim the four KPIs in 15 seconds. Land hard on EUR 49 / year — that's the price-of-coffee anchor. Then 30 seconds on the two columns: "for your company" (use the tool) vs "for your firm" (embed/refer/white-label). The split sets up slide 12.
""".strip()

NOTES_S4 = """
TIMING: ~75s. The credibility slide. Different from the banker version: here, credibility is for the BUYER — proving the math is auditable, not defending against an internal model dispute.

VOCABULARY YOU'LL NEED:
  • PD (Probability of Default) = the chance the borrower defaults within one year. Comes from S&P / Moody's long-run averages mapped to ratings. Baa1 ≈ 0.10%.
  • Risk weight K = the % of your exposure that the bank must hold as regulatory capital. 3-6% for investment grade.
  • LGD (Loss Given Default) = the fraction of exposure the bank actually loses if you default, after recovery. 30-45% typical for unsecured corporate.
  • Output floor = a Basel IV mechanism (2026: 55%, ramping to 72.5% by 2030) that prevents IRB banks from understating capital requirements. Pin from below.
  • EAD (Exposure at Default) = your committed line size adjusted for likely drawdown if you default. CCF (credit conversion factor) is 75% on the undrawn portion of confirmed RCFs.
  • CCF = Credit Conversion Factor. The fraction of an undrawn line that gets included in EAD.
  • IRB / SA = Internal Ratings-Based vs Standardised. IRB = bank's own models; SA = regulatory defaults. The output floor is "55% of SA equivalent".

LINE TO LAND ON (the highlighted row):
  "Every parameter in our model comes from your bank's own annual report. Your RM cannot say 'your tool is wrong' without first saying 'our published numbers are wrong'. They won't."
""".strip()

NOTES_S5 = """
TIMING: ~45s. Set the scene. Don't read the cards — just point.

VOCABULARY:
  • Mittelstand = German term for medium-sized industrial. Family / founder owned, EUR 100M-1B revenue, often global niche leaders. The economy's backbone.
  • BBB+ / Baa1 = the same investment-grade credit rating in S&P / Moody's notation. Implied 1-year PD ~0.10%. Typical solid mid-cap.
  • Net leverage = (Net debt) / EBITDA. 2.1x is conservative.
  • RCF = Revolving Credit Facility. Committed line you can draw, repay, redraw at will. Most common corporate credit product.
  • Bullet = principal repaid in one shot at maturity (vs. amortising = scheduled paydowns).
  • GRR = Global Guarantee Recovery Rate. Fraction of exposure covered by a third party (here: parent company guarantees 50%).
  • Confirmed = bank is contractually obligated to fund. CCF is 75% of undrawn under Basel III.

THE NARRATIVE HOOK (memorize):
  "Nordwind's CFO, Sabine, got a 150 bp quote from her BNP RM last quarter. She said yes, signed it. Three weeks later her advisor showed her this tool. Today's the post-mortem on what she could have asked for. And what she's going to ask for next time."
""".strip()

NOTES_S6 = """
TIMING: ~120s. Switch to the live tool here if possible. Pre-fill the deal (see SCENARIO.md for inputs).

VOCABULARY YOU'LL NEED LIVE:
  • Revenue = spread × drawn + commitment fee × undrawn + fees. EUR 605k/year on this deal.
  • Cost = the bank's operating cost. ~40% of revenue for plain credit products.
  • EAD = drawn + 75% × undrawn = 35M + 11.25M = 46.25M of effective exposure.
  • PD adjusted by GRR = 0.10% × (1 − 50%) = 0.05% probability of loss.
  • K = 4.26% — capital % required against this exposure. Output floor (55% of SA) is what's actually binding here.
  • FPE (Fonds Propres Économiques) = EAD × K = EUR 1.97M of BNP's equity locked up against this single facility for the life of the deal. French term, used in this engine because the original model was French. = "economic capital".
  • EL (Expected Loss) = EAD × adjusted PD = EUR 23k/year of probable losses. Tiny.
  • Net margin = Rev − Cost − Funding − EL + (Risk-free × FPE).
  • RAROC = (1 − tax) × [(Rev − Cost − Funding − EL) / FPE + Risk-free rate]. Result: 12.52%.
  • Hurdle rate = the bank's internal target. 12% is typical for European universal banks. 14-15% for US. 9% for retail-heavy.

THE ANCHOR (memorize):
  "BNP is comfortable. RAROC of 12.52% sits right above the typical 12% hurdle. The number Sabine never knew: BNP has roughly EUR 23,000 of margin per year they could give back and still hit hurdle. They were never going to volunteer that."
""".strip()

NOTES_S7 = """
TIMING: ~75s. The challenge slide. Walk the right column slowly — these are the four lines Sabine should put in her next email to the RM.

VOCABULARY:
  • Output floor binding = the IRB-computed K is below the floor, so the customer pays for the higher floor instead. A bank that doesn't yet apply 55% will look cheaper today and re-quote in 2027.
  • Internal hurdle = 12% is industry-typical but not universal. Some banks publicly target 11%, some 13%. ROTE in the annual report is the proxy.
  • Cross-sell logic = the RM may agree to lower stand-alone credit RAROC if the wallet (FX, cash, hedging, advisory) makes it up. That's leverage.

DELIVERY:
  This is the "what you walk out with" slide. The audience will lean forward — these are the four sentences they don't currently have. Read them out loud one by one. Pause after each.

DON'T ad-lib past these four — they're calibrated. Adding a fifth weakens the others.
""".strip()

NOTES_S8 = """
TIMING: ~120s. Switch to the live tool: Bank Comparison tab. Tick all 8, run.

VOCABULARY:
  • Min spread = the lowest spread that hits each bank's 12% hurdle, given everything else fixed.
  • Cheapest vs most expensive = JP Morgan can do this at 123 bp; Deutsche needs 156 bp. Same risk profile to Nordwind, 33 bp gap.

THE STORY (3 levels, memorize order):
  1. "JP Morgan would do this deal at 123 bp. BNP needs 143. Same risk to Nordwind. That's a 20 basis point gap that's now visible to you."
  2. "Deutsche Bank can't do it under 156 — they're losing money below that. So when Deutsche says 'best we can do is 145', they're being generous, not stingy. Knowing the difference matters."
  3. "Nordwind doesn't have to switch banks. They have to walk into the room with the chart on slide 8."

TONE WARNING:
  Don't trash any specific bank by name. JPM's high-RAROC isn't them being cheap; it's tax + scale. Deutsche's high min-spread isn't incompetence; it's CTI + funding cost. Frame the dispersion as structural, not malicious.
""".strip()

NOTES_S9 = """
TIMING: ~75s. Three columns, ~25s each. Land on CTI hardest — it's the most defensible point.

VOCABULARY:
  • CTI = Cost-to-Income ratio. Operating costs / net banking income. Public, audited, in every annual report. Bank's headline efficiency number. Lower = leaner.
  • Effective tax rate (ETR) = actual tax paid / pre-tax income. Geography-driven mostly. US banks ~21% post-TCJA, EU 25-28%.
  • Output floor / LGD floors = Basel IV mechanisms. LGD floors: 25% unsecured, 10% receivables-secured, 0% financial collateral. Phase-in 2025-2030.

THE PIVOT TO THE CFO ACTION (memorize):
  "You're not negotiating the math. You're using the math to negotiate. Walk into the meeting saying 'your CTI is 62% — ten points above ING — and that's making my deal more expensive than it should be'. The RM cannot answer that with a vague 'market conditions' anymore."
""".strip()

NOTES_S10 = """
TIMING: ~90s. Live: same deal in the tool, click the "Solve minimum spread" button, target = 12%, returns 143 bp.

VOCABULARY:
  • Brent's method = numerical root finder. 5 ms per call. Don't go deep unless asked.
  • "Run backwards" = same engine, RAROC as input, spread as output.

THE PUNCHLINE (memorize):
  "Sabine walks into the next meeting with this table on her iPad. Her opening line is 'I know your floor on this deal is 143. Let's start there'. The negotiation that used to take three calls now takes thirty seconds. The RM either matches, walks, or counters with a bundle. All three outcomes are wins for Sabine compared to where she started."

DELIVERY HACK:
  If you have a tablet or second laptop, actually mirror the live tool while you say this. The visual of the table appearing on demand sells the moment.
""".strip()

NOTES_S11 = """
TIMING: ~75s. The outcome slide. Three plausible endings. Walk left to right. Frame as "the spectrum Sabine controls now".

THE THREE OUTCOMES IN ONE LINE EACH:
  1. RM matches at 143 → Sabine saves EUR 23k / year × 5 years = EUR 116k locked in.
  2. RM bundles down to 135 + cross-sell concession → Sabine saves ~EUR 75k / year and locks in FX pricing for 3 years. Best outcome.
  3. RM holds at 150 → Sabine takes the file to ING (12.25% RAROC at 150 bp, has more room). BNP loses the renewal.

ALL THREE ARE WINS FOR SABINE. That's the killer point. Today she has zero of these levers. After OpenRAROC she has all three.

THE LINE TO LAND ON:
  "The deal still gets done. The question is whether you're driving the conversation or being driven by it."
""".strip()

NOTES_S12 = """
TIMING: ~75s. Two doors. One for each side of the room. Don't pitch — just open the doors.

DOOR A (corporate audience):
  • Free tier: 4 banks, full calculator, full bank comparison. Plenty for one-off use on a single facility.
  • Pro: EUR 49/year, all 59 banks, portfolio optimizer (allocate facilities across banks to minimise total cost subject to wallet constraints), PDF reports, API access.
  • Onboarding: 30 minutes from signup to your first portfolio comparison. No procurement cycle at this price point.

DOOR B (partner / advisor audience):
  • Referral: send your clients to openraroc.com, get 20% recurring. White-glove onboarding for your top tier.
  • White-label: brand the engine + bank profiles as your own product. Same calculator, your logo, your domain.
  • OEM: deeper embed via API + MCP server. Plug RAROC into your existing client portal, advisory dashboard, or treasury tool.

CLOSING LINE (memorize, do not paraphrase):
  "Two days from now, half of you will have run a deal through the free tier on a real facility. One in ten of you will be on Pro within a month. Three of you will email me about the partner program. That's the predictable distribution. The unpredictable part is which three."

THEN: pause. Take questions. Have the laptop and openraroc.com open behind you for the live walkthrough requests.
""".strip()


# ───────────────── Slide 1 — Title ─────────────────────────────────
s = new_slide(page=None, notes=NOTES_S1)
band = add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5),
                fill=ACCENT, line=None)
add_text(s, Inches(0.7), Inches(0.6), Inches(8), Inches(0.5),
         "OpenRAROC", size=22, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.1), Inches(8), Inches(0.4),
         "openraroc.com", size=12, color=ACCENT)

add_text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(1.5),
         "Negotiate your bank pricing\nwith the math your bank uses internally.",
         size=42, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.7),
         "A live demo — for treasurers and the advisors who serve them.",
         size=22, color=ACCENT)

add_line(s, Inches(0.7), Inches(5.4), Inches(5), Inches(5.4), color=BORDER, width=1.0)
add_text(s, Inches(0.7), Inches(5.55), Inches(8), Inches(0.4),
         "10–15 minutes  •  one realistic mid-cap deal, three chapters",
         size=14, color=TEXT2)
add_text(s, Inches(0.7), Inches(5.95), Inches(8), Inches(0.4),
         "Live walk-through in the tool — bring questions",
         size=12, color=TEXT3)

add_text(s, Inches(0.7), Inches(7.0), Inches(8), Inches(0.3),
         "Cyril Poder  •  cyril.poder@gmail.com", size=10, color=TEXT3)
add_text(s, Inches(11.0), Inches(7.0), Inches(2), Inches(0.3),
         "1 / 12", size=10, color=TEXT3, align=PP_ALIGN.RIGHT)


# ───────────────── Slide 2 — The opacity problem ──────────────────
s = new_slide("THE PROBLEM",
              "Your bank quotes one number. They model fifteen.",
              "And you negotiate against the one.", page=2, notes=NOTES_S2)

add_bullet_block(s, Inches(0.5), Inches(2.1), Inches(6.1), Inches(2.2),
    "What you see on the term sheet",
    [
        "\"150 bp over EURIBOR + 20 bp commitment fee.\"",
        "\"Subject to credit committee approval.\"",
        "\"Best we can do given current market conditions.\"",
        "Total information you have to push back: zero.",
    ],
    title_color=TEXT2, fill=SURFACE)

add_bullet_block(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(2.2),
    "What your RM defends internally",
    [
        "Risk weight K, FPE, EAD, expected loss.",
        "Funding cost, output floor, hurdle rate.",
        "Cost-to-income ratio, effective tax rate.",
        "All public. All extractable. All in OpenRAROC.",
    ],
    title_color=ACCENT, fill=SURFACE)

add_rect(s, Inches(0.5), Inches(4.55), Inches(12.4), Inches(2.0),
         fill=SURFACE, line=ACCENT, line_w=1.5)
add_text(s, Inches(0.8), Inches(4.78), Inches(12), Inches(0.5),
         "Why this matters in money", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(5.13), Inches(12), Inches(1.5),
         "30 basis points on a EUR 50 million facility is EUR 750,000 over a 5-year tenor. "
         "That's a senior hire. That's a software stack. That's the dividend your CEO didn't get to announce. "
         "And you negotiated it away in a 45-minute call because you didn't have the math.",
         size=14, color=TEXT, anchor=MSO_ANCHOR.TOP)


# ───────────────── Slide 3 — What we built ─────────────────────────
s = new_slide("THE TOOL", "What OpenRAROC is.",
              "Basel III IRB calculator + 59 bank profiles + min-spread solver.",
              page=3, notes=NOTES_S3)

add_kpi(s, Inches(0.5), Inches(2.1), Inches(2.95), Inches(1.4),
        "Bank profiles", "59",
        "from public Pillar 3 disclosures", value_color=ACCENT)
add_kpi(s, Inches(3.6), Inches(2.1), Inches(2.95), Inches(1.4),
        "Countries", "26",
        "EU + UK + US + APAC + LatAm + Gulf", value_color=ACCENT)
add_kpi(s, Inches(6.7), Inches(2.1), Inches(2.95), Inches(1.4),
        "Products", "7",
        "Term, RCF, LC, Bond, IRS, FX, Project finance",
        value_color=ACCENT)
add_kpi(s, Inches(9.8), Inches(2.1), Inches(3.05), Inches(1.4),
        "Pro tier", "EUR 49 / yr",
        "4 banks free, 55 unlocked at Pro", value_color=ACCENT)

add_bullet_block(s, Inches(0.5), Inches(3.7), Inches(6.1), Inches(3.0),
    "For your company",
    [
        "Run any facility through the calculator in under 60 seconds.",
        "Compare your current bank against 58 alternatives.",
        "Reverse-solve the minimum spread your bank can accept.",
        "Generate a CFO-ready PDF for the credit committee.",
        "Set up a portfolio across all your facilities and rebalance.",
    ], title_color=ACCENT)

add_bullet_block(s, Inches(6.85), Inches(3.7), Inches(6.0), Inches(3.0),
    "For your firm  (advisors, consultancies, fractional CFOs)",
    [
        "Embed the engine in your client deliverables.",
        "Three flavours: referral (recurring), white-label, full OEM.",
        "API + MCP — wire RAROC into your treasury or advisory tool.",
        "Annual data refresh; no maintenance burden on your side.",
        "Defensible, auditable, regulator-aligned methodology.",
    ], title_color=GREEN)


# ───────────────── Slide 4 — Methodology ────────────────────────────
s = new_slide("THE TOOL", "Why your bank cannot dispute the math.",
              "Three public sources. Zero estimates. Auditable end to end.",
              page=4, notes=NOTES_S4)

data = [
    ["Input", "Where it comes from", "What it drives"],
    ["Probability of Default (PD)", "S&P Global long-run corporate default rates",
     "Rating → PD lookup. Floored at 5 bp under Basel III."],
    ["Risk weight K, b, R", "BIS CRE31 / CRE32 (the Basel framework itself)",
     "Capital requirement formula — same across all IRB banks."],
    ["LGD, CTI, tax, funding spread", "Each bank's annual Pillar 3 + 10-K / URD",
     "Bank-specific economics that move RAROC."],
    ["Output floor (Basel IV)", "EBA / BIS phase-in schedule",
     "55% in 2026 → 72.5% in 2030. Floors capital downside."],
    ["Credit Conversion Factors", "BCBS / EBA confirmed vs unconfirmed CCFs",
     "Drives EAD on the undrawn portion of revolvers."],
]
add_table(s, Inches(0.5), Inches(2.1), Inches(12.4), Inches(3.6), data,
          col_widths=[Inches(2.6), Inches(4.5), Inches(5.3)],
          highlight_rows=[2], highlight_color=ACCENT)

add_rect(s, Inches(0.5), Inches(5.85), Inches(12.4), Inches(0.95),
         fill=SURFACE, line=BORDER)
add_text(s, Inches(0.75), Inches(6.0), Inches(12), Inches(0.35),
         "Why this matters in your next meeting",
         size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.75), Inches(6.32), Inches(12), Inches(0.5),
         "Every input in our model comes from your bank's own annual report. "
         "Your RM cannot say \"your tool is wrong\" without first saying \"our published filings are wrong\". They won't.",
         size=11, color=TEXT2)


# ───────────────── Slide 5 — Scenario intro ─────────────────────────
s = new_slide("THE STORY", "Meet Sabine. She's the CFO of Nordwind.",
              "She just signed the wrong deal. Today's the post-mortem — and the playbook for next time.",
              page=5, notes=NOTES_S5)

add_rect(s, Inches(0.5), Inches(2.1), Inches(5.6), Inches(4.65),
         fill=SURFACE, line=BORDER)
add_text(s, Inches(0.75), Inches(2.3), Inches(5), Inches(0.4),
         "THE COMPANY", size=10, bold=True, color=TEXT3)
add_text(s, Inches(0.75), Inches(2.6), Inches(5), Inches(0.5),
         "Nordwind Industries GmbH", size=20, bold=True, color=WHITE)
add_text(s, Inches(0.75), Inches(3.1), Inches(5), Inches(0.4),
         "German Mittelstand • Industrial machinery", size=12, color=ACCENT)

facts = [
    ("Revenue (FY)", "EUR 820 m"),
    ("EBITDA margin", "11.4 %"),
    ("Net leverage", "2.1x"),
    ("Public rating", "BBB+ / Baa1"),
    ("Parent guarantee", "Yes — covers 50% (GRR)"),
    ("Existing relationships", "BNP, HSBC, Deutsche, ING, SocGen"),
]
y0 = 3.65
for i, (k, v) in enumerate(facts):
    add_text(s, Inches(0.75), Inches(y0 + i*0.42), Inches(2.8), Inches(0.35),
             k, size=11, color=TEXT2)
    add_text(s, Inches(3.55), Inches(y0 + i*0.42), Inches(2.5), Inches(0.35),
             v, size=11, bold=True, color=TEXT)

add_rect(s, Inches(6.35), Inches(2.1), Inches(6.5), Inches(4.65),
         fill=SURFACE, line=ACCENT, line_w=1.5)
add_text(s, Inches(6.6), Inches(2.3), Inches(6), Inches(0.4),
         "THE DEAL SHE SIGNED LAST QUARTER",
         size=10, bold=True, color=ACCENT)
add_text(s, Inches(6.6), Inches(2.6), Inches(6), Inches(0.5),
         "5-year revolving credit facility", size=20, bold=True, color=WHITE)
add_text(s, Inches(6.6), Inches(3.1), Inches(6), Inches(0.4),
         "Refi + general corporate purposes — quoted by BNP",
         size=12, color=TEXT2)

deal_facts = [
    ("Committed amount", "EUR 50 m"),
    ("Expected average drawn", "EUR 35 m  (70 %)"),
    ("Tenor", "5 years (60 m), bullet"),
    ("Spread", "150 bp"),
    ("Commitment fee (undrawn)", "20 bp"),
    ("Participation fee (one-off)", "EUR 50 k"),
    ("Confirmed / Committed", "Yes"),
]
y0 = 3.6
for i, (k, v) in enumerate(deal_facts):
    add_text(s, Inches(6.6), Inches(y0 + i*0.41), Inches(3), Inches(0.35),
             k, size=11, color=TEXT2)
    add_text(s, Inches(9.6), Inches(y0 + i*0.41), Inches(3.2), Inches(0.35),
             v, size=11, bold=True, color=TEXT)


# ───────────────── Slide 6 — Chapter 1 deep-dive ───────────────────
s = new_slide("CHAPTER 1 — WHAT YOUR BANK SEES",
              "BNP's RAROC on Sabine's deal, line by line.",
              "Same numbers BNP's pricing committee approved against. Now visible to you.",
              page=6, notes=NOTES_S6)

data = [
    ["Line", "Value", "Plain English"],
    ["Revenue",        "EUR 605,000",   "What BNP earns annually from spread + fees"],
    ["Operating cost", "EUR 242,000",   "BNP's cost to run the relationship (~40% of rev)"],
    ["EAD",            "EUR 46,250,000","BNP's exposure if you default (drawn + 75% undrawn)"],
    ["PD (Baa1)",      "0.10 %",        "Annual probability of default at your rating"],
    ["PD adj. for GRR","0.05 %",        "After the parent guarantee absorbs half the loss"],
    ["Risk weight K",  "4.26 %",        "% of EAD BNP holds as regulatory capital"],
    ["FPE (capital)",  "EUR 1,971,714", "BNP's equity locked up against your facility"],
    ["Expected loss",  "EUR 23,125",    "Average annual loss across many possible futures"],
    ["Net margin",     "EUR 334,581",   "Pre-tax economic profit on the deal"],
    ["Tax (BNP, 26%)", "EUR 87,660",    "French effective tax rate from BNP's URD 2024"],
]
add_table(s, Inches(0.5), Inches(2.0), Inches(8.0), Inches(4.5), data,
          col_widths=[Inches(2.0), Inches(2.4), Inches(3.6)],
          font_size=10, header_size=10,
          num_cols_right=[1])

add_rect(s, Inches(8.7), Inches(2.0), Inches(4.2), Inches(4.5),
         fill=SURFACE, line=ACCENT, line_w=1.5)
add_text(s, Inches(8.9), Inches(2.2), Inches(4), Inches(0.4),
         "BNP's RAROC", size=10, bold=True, color=ACCENT)
add_text(s, Inches(8.9), Inches(2.55), Inches(4), Inches(0.7),
         "Computed", size=14, color=TEXT2)
add_text(s, Inches(8.9), Inches(2.95), Inches(4), Inches(1.3),
         "12.52 %", size=54, bold=True, color=GREEN)
add_text(s, Inches(8.9), Inches(4.3), Inches(4), Inches(0.4),
         "vs ~12 % typical European hurdle",
         size=11, color=TEXT2)

add_line(s, Inches(8.9), Inches(4.85), Inches(12.7), Inches(4.85),
         color=BORDER, width=0.75)

add_text(s, Inches(8.9), Inches(4.95), Inches(4), Inches(0.35),
         "WHAT THIS MEANS FOR SABINE",
         size=9, bold=True, color=TEXT3)
add_text(s, Inches(8.9), Inches(5.25), Inches(4), Inches(0.4),
         "EUR 23k / year of comfort",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(8.9), Inches(5.7), Inches(4), Inches(0.5),
         "BNP can give that back and still hit hurdle.\nThey were never going to volunteer it.",
         size=11, color=TEXT2)

add_text(s, Inches(0.5), Inches(6.65), Inches(12), Inches(0.4),
         "Live walk-through — type any deal, every line moves in real time.",
         size=11, color=TEXT3)


# ───────────────── Slide 7 — Chapter 1 commentary ──────────────────
s = new_slide("CHAPTER 1 — WHAT YOUR BANK SEES",
              "Four sentences Sabine couldn't say last quarter.",
              "Each one is a lever. Each one is now in her email draft.",
              page=7, notes=NOTES_S7)

add_bullet_block(s, Inches(0.5), Inches(2.1), Inches(12.4), Inches(4.5),
    "What Sabine sends to her RM next week",
    [
        "\"Your RAROC on this deal is 12.5%. Your group hurdle is 12. We agreed at the top of your range — let's revisit.\"",
        "",
        "\"You priced the parent guarantee in the LGD, which is correct. But did you also adjust the maturity-band? My read is no.\"",
        "",
        "\"I see your Q3 disclosure has the output floor at 55%. If your internal model has not yet caught up, my pricing is structurally too high.\"",
        "",
        "\"BNP CTI is 62%. ING runs 55%. I am paying for ten points of inefficiency that has nothing to do with my credit.\"",
    ], title_color=ACCENT)


# ───────────────── Slide 8 — Chapter 2 comparison ──────────────────
s = new_slide("CHAPTER 2 — HOW YOUR BANK RANKS",
              "Same deal. Eight banks. Five of them already lend to Nordwind.",
              "Real numbers — engine output as of today.", page=8, notes=NOTES_S8)

data = [
    ["Bank", "Country", "CTI", "Tax", "RAROC at 150 bp", "Min spread to hit 12%"],
    ["JP Morgan",       "United States",  "52 %", "21 %", "14.26 %", "123 bp"],
    ["Crédit Agricole", "France",         "56 %", "24 %", "13.43 %", "132 bp"],
    ["Barclays",        "United Kingdom", "62 %", "22 %", "13.30 %", "134 bp"],
    ["HSBC",            "United Kingdom", "53 %", "23 %", "13.12 %", "136 bp"],
    ["Société Générale","France",         "64 %", "20 %", "12.62 %", "143 bp"],
    ["BNP Paribas",     "France",         "62 %", "26 %", "12.52 %", "143 bp"],
    ["ING Group",       "Netherlands",    "55 %", "28 %", "12.25 %", "147 bp"],
    ["Deutsche Bank",   "Germany",        "64 %", "27 %", "11.56 %", "156 bp"],
]
add_table(s, Inches(0.5), Inches(2.1), Inches(12.4), Inches(4.0), data,
          col_widths=[Inches(2.4), Inches(2.0), Inches(1.2), Inches(1.2),
                      Inches(2.8), Inches(2.8)],
          highlight_rows=[5], highlight_color=ACCENT,
          font_size=11, header_size=10,
          num_cols_right=[2, 3, 4, 5])

add_text(s, Inches(0.5), Inches(6.25), Inches(12), Inches(0.4),
         "33 bp gap between cheapest and most expensive — same risk to Nordwind.",
         size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.5), Inches(6.65), Inches(12), Inches(0.4),
         "Sabine doesn't need to switch banks. She needs to walk in with this chart.",
         size=12, color=TEXT2)


# ───────────────── Slide 9 — Why the spread ─────────────────────────
s = new_slide("CHAPTER 2 — HOW YOUR BANK RANKS",
              "Three structural reasons your bank is more expensive than someone else's.",
              "All three are now things you can name out loud in the meeting.",
              page=9, notes=NOTES_S9)

s_ = s
def slide_text(x, y, w, h, paragraphs):
    tb = s_.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, p_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        p.space_after = Pt(8)
        r = p.add_run(); r.text = p_text
        r.font.size = Pt(11); r.font.color.rgb = TEXT
        r.font.name = "Calibri"

def col(x, title, body, accent_color):
    add_rect(s, x, Inches(2.1), Inches(4.0), Inches(4.6),
             fill=SURFACE, line=accent_color, line_w=1.25)
    add_text(s, x + Inches(0.25), Inches(2.3), Inches(3.6), Inches(0.4),
             title, size=14, bold=True, color=accent_color)
    slide_text(x + Inches(0.25), Inches(2.7), Inches(3.6), Inches(3.9), body)

col(Inches(0.5), "Cost-to-income (CTI)",
    [
      "Your bank spends 64 cents to earn each euro of revenue. Their best-in-class peer spends 52 cents.",
      "The 12-cent gap shows up in your spread.",
      "Public, audited, on page 1 of every annual report. Your RM cannot dispute this number.",
      "Your line: \"I'm paying for your inefficiency.\"",
    ], ORANGE)

col(Inches(4.65), "Effective tax rate",
    [
      "RAROC is post-tax. A US bank pays 21% effective; a Dutch bank pays 28%.",
      "5 points of tax delta = ~70 bp of RAROC delta on this deal.",
      "Pure geography. The bank with the better passport quotes you better.",
      "Your line: \"Have you considered which entity books this deal?\"",
    ], ACCENT)

col(Inches(8.8), "Output floor & LGD floor",
    [
      "Banks not yet at 55% output floor look cheaper today and re-price you in 2027.",
      "LGD floors (25% unsecured / 10% receivables-secured) hit each bank differently.",
      "Both are Basel IV mechanisms phasing in 2025-2030.",
      "Your line: \"What floor are you currently applying to my exposure?\"",
    ], GREEN)


# ───────────────── Slide 10 — Min-spread solver ─────────────────────
s = new_slide("CHAPTER 3 — THE NEGOTIATION",
              "What Sabine brings to her next meeting with BNP.",
              "The same engine, run backwards. Asks: what's the minimum spread BNP can accept?",
              page=10, notes=NOTES_S10)

add_rect(s, Inches(0.5), Inches(2.1), Inches(5.5), Inches(4.5),
         fill=SURFACE, line=BORDER)
add_text(s, Inches(0.75), Inches(2.3), Inches(5), Inches(0.4),
         "THE QUESTION SHE NOW CAN ASK", size=10, bold=True, color=ACCENT)
add_text(s, Inches(0.75), Inches(2.65), Inches(5), Inches(0.6),
         "\"What is the lowest spread\nyour committee will approve?\"",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(0.75), Inches(3.85), Inches(5), Inches(0.4),
         "WHAT THE TOOL DOES", size=10, bold=True, color=ACCENT)
add_text(s, Inches(0.75), Inches(4.2), Inches(5), Inches(2.0),
         "Holds everything else constant.\n"
         "Solves for the spread that hits 12%.\n\n"
         "Result for BNP on Sabine's deal: 143 bp.\n"
         "20 basis points below the 150 they quoted.",
         size=12, color=TEXT2)

data = [
    ["Spread", "RAROC", "Annual revenue", "vs current"],
    ["150 bp",  "12.52 %", "EUR 605,000", "—"],
    ["143 bp",  "12.00 %", "EUR 581,705", "−EUR 23,295"],
    ["140 bp",  "11.74 %", "EUR 570,000", "−EUR 35,000"],
    ["130 bp",  "10.95 %", "EUR 535,000", "−EUR 70,000"],
    ["120 bp",  "10.17 %", "EUR 500,000", "−EUR 105,000"],
]
add_table(s, Inches(6.25), Inches(2.1), Inches(6.6), Inches(4.0), data,
          col_widths=[Inches(1.4), Inches(1.5), Inches(2.0), Inches(1.7)],
          highlight_rows=[1], highlight_color=ACCENT,
          font_size=12, header_size=10,
          num_cols_right=[0, 1, 2, 3])

add_text(s, Inches(6.25), Inches(6.25), Inches(6.6), Inches(0.4),
         "Sabine walks in saying: \"I know your floor is 143. Let's start there.\"",
         size=12, bold=True, color=ACCENT)
add_text(s, Inches(6.25), Inches(6.6), Inches(6.6), Inches(0.4),
         "The negotiation that used to take three calls now takes thirty seconds.",
         size=11, color=TEXT2)


# ───────────────── Slide 11 — Negotiation outcomes ──────────────────
s = new_slide("CHAPTER 3 — THE NEGOTIATION",
              "Three plausible endings. Sabine wins all three.",
              "Today she has none of these levers. After OpenRAROC she has all of them.",
              page=11, notes=NOTES_S11)

scenarios = [
    ("BNP matches at 143",
     ACCENT,
     ["Spread cut from 150 → 143 bp.",
      "Saving: EUR 23k / year × 5 yrs = EUR 116k locked in.",
      "Relationship intact. Cleanest outcome.",
     ]),
    ("BNP bundles at 135",
     GREEN,
     ["RM offers 135 bp + concession on FX & cash mandate.",
      "Wallet RAROC stays > 12% across the relationship.",
      "Saving: ~EUR 75k / year + locked-in FX pricing. Best case.",
     ]),
    ("BNP holds at 150",
     YELLOW,
     ["BNP refuses. Sabine takes the file to ING (12.25% at 150).",
      "BNP loses the renewal but keeps cash management.",
      "Sabine still wins — better pricing at a different bank.",
     ]),
]
x_positions = [Inches(0.5), Inches(4.78), Inches(9.05)]
for x, (title, color, bullets) in zip(x_positions, scenarios):
    add_rect(s, x, Inches(2.1), Inches(4.0), Inches(4.5),
             fill=SURFACE, line=color, line_w=1.5)
    add_text(s, x + Inches(0.25), Inches(2.3), Inches(3.6), Inches(0.4),
             title, size=13, bold=True, color=color)
    tb = s.shapes.add_textbox(x + Inches(0.25), Inches(2.85),
                              Inches(3.6), Inches(3.6))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3; p.space_after = Pt(8)
        r = p.add_run(); r.text = "•  " + b
        r.font.size = Pt(11); r.font.color.rgb = TEXT
        r.font.name = "Calibri"

add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
         "The deal still gets done. The question is whether you're driving the conversation or being driven by it.",
         size=12, color=TEXT2, align=PP_ALIGN.CENTER)


# ───────────────── Slide 12 — Two doors ──────────────────────────────
s = new_slide("WHAT'S NEXT", "Two doors. Pick the one for your role.",
              "Both open today. No procurement cycle.",
              page=12, notes=NOTES_S12)

add_rect(s, Inches(0.5), Inches(2.1), Inches(6.1), Inches(4.55),
         fill=SURFACE, line=ACCENT, line_w=1.5)
add_text(s, Inches(0.75), Inches(2.3), Inches(5.5), Inches(0.4),
         "FOR YOUR COMPANY", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.75), Inches(2.65), Inches(5.5), Inches(0.6),
         "Use the tool", size=22, bold=True, color=WHITE)
add_text(s, Inches(0.75), Inches(3.25), Inches(5.5), Inches(0.4),
         "Treasurers, CFOs, finance directors", size=11, color=TEXT2)

corp_bullets = [
    ("Free tier",
     "4 banks. Full calculator, comparison, sensitivity. Enough for a single facility."),
    ("Pro — EUR 49 / year",
     "All 59 banks. Portfolio optimizer. PDF reports for the credit committee. API + CLI access."),
    ("Onboarding",
     "30 minutes from sign-up to your first bank comparison. No procurement cycle at this price."),
    ("Try right now",
     "openraroc.com — runs on your laptop, no install, no card on file for the free tier."),
]
y0 = 3.85
for label, body in corp_bullets:
    add_text(s, Inches(0.75), Inches(y0), Inches(5.5), Inches(0.3),
             label, size=11, bold=True, color=ACCENT)
    add_text(s, Inches(0.75), Inches(y0 + 0.28), Inches(5.5), Inches(0.4),
             body, size=10, color=TEXT2)
    y0 += 0.68

add_rect(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(4.55),
         fill=SURFACE, line=GREEN, line_w=1.5)
add_text(s, Inches(7.1), Inches(2.3), Inches(5.5), Inches(0.4),
         "FOR YOUR FIRM", size=11, bold=True, color=GREEN)
add_text(s, Inches(7.1), Inches(2.65), Inches(5.5), Inches(0.6),
         "Bring it to your clients", size=22, bold=True, color=WHITE)
add_text(s, Inches(7.1), Inches(3.25), Inches(5.5), Inches(0.4),
         "Treasury consultancies, debt advisors, fractional CFOs, accountants",
         size=11, color=TEXT2)

partner_bullets = [
    ("Referral",
     "Send clients to openraroc.com with your tracking link. 20% recurring."),
    ("White-label",
     "Same engine and bank profiles, your brand, your domain. Annual licence."),
    ("OEM / embed",
     "Plug RAROC into your existing client portal via API + MCP server."),
    ("All three",
     "Start as a referrer this quarter; upgrade to white-label when you have 5+ active clients."),
]
y0 = 3.85
for label, body in partner_bullets:
    add_text(s, Inches(7.1), Inches(y0), Inches(5.5), Inches(0.3),
             label, size=11, bold=True, color=GREEN)
    add_text(s, Inches(7.1), Inches(y0 + 0.28), Inches(5.5), Inches(0.4),
             body, size=10, color=TEXT2)
    y0 += 0.68

add_line(s, Inches(0.5), Inches(6.78), Inches(12.83), Inches(6.78),
         color=BORDER, width=0.75)
add_text(s, Inches(0.5), Inches(6.85), Inches(8), Inches(0.3),
         "Cyril Poder  •  cyril.poder@gmail.com  •  openraroc.com",
         size=11, color=TEXT2)
add_text(s, Inches(11), Inches(6.85), Inches(2), Inches(0.3),
         "Q & A", size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


# ── Save ──────────────────────────────────────────────────────────────
out = "/home/cpo/raroc/demo_deck/openraroc_demo.pptx"
prs.save(out)
print(f"Saved: {out}")
